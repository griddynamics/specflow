"""Tests for services.document_reader.

Each test creates minimal real files via the respective library so
extraction runs against actual format internals, not mocks.
"""

from __future__ import annotations

import base64
import csv
import io
import tempfile
from pathlib import Path

import pytest
from docx import Document as DocxDocument
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from services.document_reader import (
    DocumentContent,
    EmbeddedImage,
    is_image_file,
    is_supported_document,
    read_document,
)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

# Minimal valid PDF bytes (written to tmp_path at test time; not stored as files in repo).
_TWO_PAGE_TEXT_PDF = base64.b64decode(
    "JVBERi0xLjMKJZOMi54gUmVwb3J0TGFiIEdlbmVyYXRlZCBQREYgZG9jdW1lbnQgKG9wZW5zb3VyY2UpCjEgMCBvYmoK"
    "PDwKL0YxIDIgMCBSCj4+CmVuZG9iagoyIDAgb2JqCjw8Ci9CYXNlRm9udCAvSGVsdmV0aWNhIC9FbmNvZGluZyAv"
    "V2luQW5zaUVuY29kaW5nIC9OYW1lIC9GMSAvU3VidHlwZSAvVHlwZTEgL1R5cGUgL0ZvbnQKPj4KZW5kb2JqCjMg"
    "MCBvYmoKPDwKL0NvbnRlbnRzIDggMCBSIC9NZWRpYUJveCBbIDAgMCA2MTIgNzkyIF0gL1BhcmVudCA3IDAgUiAv"
    "UmVzb3VyY2VzIDw8Ci9Gb250IDEgMCBSIC9Qcm9jU2V0IFsgL1BERiAvVGV4dCAvSW1hZ2VCIC9JbWFnZUMgL0lt"
    "YWdlSSBdCj4+IC9Sb3RhdGUgMCAvVHJhbnMgPDwKCj4+IAogIC9UeXBlIC9QYWdlCj4+CmVuZG9iago0IDAgb2Jq"
    "Cjw8Ci9Db250ZW50cyA5IDAgUiAvTWVkaWFCb3ggWyAwIDAgNjEyIDc5MiBdIC9QYXJlbnQgNyAwIFIgL1Jlc291"
    "cmNlcyA8PAovRm9udCAxIDAgUiAvUHJvY1NldCBbIC9QREYgL1RleHQgL0ltYWdlQiAvSW1hZ2VDIC9JbWFnZUkg"
    "XQo+PiAvUm90YXRlIDAgL1RyYW5zIDw8Cgo+PiAKICAvVHlwZSAvUGFnZQo+PgplbmRvYmoKNSAwIG9iago8PAov"
    "UGFnZU1vZGUgL1VzZU5vbmUgL1BhZ2VzIDcgMCBSIC9UeXBlIC9DYXRhbG9nCj4+CmVuZG9iago2IDAgb2JqCjw8"
    "Ci9BdXRob3IgKGFub255bW91cykgL0NyZWF0aW9uRGF0ZSAoRDoyMDI2MDYzMDEwMjE0MyswMicwMCcpIC9DcmVhdG9y"
    "IChhbm9ueW1vdXMpIC9LZXl3b3JkcyAoKSAvTW9kRGF0ZSAoRDoyMDI2MDYzMDEwMjE0MyswMicwMCcpIC9Qcm9kdWNl"
    "ciAoUmVwb3J0TGFiIFBERiBMaWJyYXJ5IC0gXChvcGVuc291cmNlXCkpIAogIC9TdWJqZWN0ICh1bnNwZWNpZmll"
    "ZCkgL1RpdGxlICh1bnRpdGxlZCkgL1RyYXBwZWQgL0ZhbHNlCj4+CmVuZG9iago3IDAgb2JqCjw8Ci9Db3VudCAy"
    "IC9LaWRzIFsgMyAwIFIgNCAwIFIgXSAvVHlwZSAvUGFnZXMKPj4KZW5kb2JqCjggMCBvYmoKPDwKL0ZpbHRlciBb"
    "IC9BU0NJSTg1RGVjb2RlIC9GbGF0ZURlY29kZSBdIC9MZW5ndGggMTA3Cj4+CnN0cmVhbQpHYXBRaDBFPUYsMFVc"
    "SDNUXHBOWVReUUtrP3RjPklQLDtXI1UxXjIzaWhQRU1fP0NXNEtJU2k8IVs3YCNPQl9zS0pBaChaTCZBU19MI2kw"
    "RVAxNyxXOW5lWltLYixodCJwSFVaR3Azc2F+PmVuZHN0cmVhbQplbmRvYmoKOSAwIG9iago8PAovRmlsdGVyIFsg"
    "L0FTQ0lJODVEZWNvZGUgL0ZsYXRlRGVjb2RlIF0gL0xlbmd0aCAxMDcKPj4Kc3RyZWFtCkdhcFFoMEU9RiwwVVxI"
    "M1RccE5ZVF5RS2s/dGM+SVAsO1cjVTFeMjNpaFBFTV8/Q1c0S0lTaTwhWzdgI09CX3NLSkFoKFpLcCkiU0wjaTBF"
    "UDE3LFc5bmVaW0tiLGh0InBIVVpIVVNZSX4+ZW5kc3RyZWFtCmVuZG9iagp4cmVmCjAgMTAKMDAwMDAwMDAwMCA2"
    "NTUzNSBmIAowMDAwMDAwMDYxIDAwMDAwIG4gCjAwMDAwMDAwOTIgMDAwMDAgbiAKMDAwMDAwMDE5OSAwMDAwMCBu"
    "IAowMDAwMDAwMzkyIDAwMDAwIG4gCjAwMDAwMDA1ODUgMDAwMDAgbiAKMDAwMDAwMDY1MyAwMDAwMCBuIAowMDAw"
    "MDAwOTE0IDAwMDAwIG4gCjAwMDAwMDA5NzkgMDAwMDAgbiAKMDAwMDAwMTE3NiAwMDAwMCBuIAp0cmFpbGVyCjw8"
    "Ci9JRCAKWzwyNWM4ZGVjMGM0ODY2NzY4ZWI5NGNkMDdlYWQ4NjIwYj48MjVjOGRlYzBjNDg2Njc2OGViOTRjZDA3"
    "ZWFkODYyMGI+XQolIFJlcG9ydExhYiBnZW5lcmF0ZWQgUERGIGRvY3VtZW50IC0tIGRpZ2VzdCAob3BlbnNvdXJj"
    "ZSkKCi9JbmZvIDYgMCBSCi9Sb290IDUgMCBSCi9TaXplIDEwCj4+CnN0YXJ0eHJlZgoxMzczCiUlRU9GCg=="
)
_PDF_WITH_EMBEDDED_IMAGE = base64.b64decode(
    "JVBERi0xLjMKJZOMi54gUmVwb3J0TGFiIEdlbmVyYXRlZCBQREYgZG9jdW1lbnQgKG9wZW5zb3VyY2UpCjEgMCBv"
    "YmoKPDwKL0YxIDIgMCBSCj4+CmVuZG9iagoyIDAgb2JqCjw8Ci9CYXNlRm9udCAvSGVsdmV0aWNhIC9FbmNvZGlu"
    "ZyAvV2luQW5zaUVuY29kaW5nIC9OYW1lIC9GMSAvU3VidHlwZSAvVHlwZTEgL1R5cGUgL0ZvbnQKPj4KZW5kb2Jq"
    "CjMgMCBvYmoKPDwKL0JpdHNQZXJDb21wb25lbnQgOCAvQ29sb3JTcGFjZSAvRGV2aWNlUkdCIC9GaWx0ZXIgWyAv"
    "QVNDSUk4NURlY29kZSAvRmxhdGVEZWNvZGUgXSAvSGVpZ2h0IDQgL0xlbmd0aCAyMCAvU3VidHlwZSAvSW1hZ2Ug"
    "CiAgL1R5cGUgL1hPYmplY3QgL1dpZHRoIDQKPj4Kc3RyZWFtCkdiIlpXXyFtczEiOT1hYyYrS34+ZW5kc3RyZWFt"
    "CmVuZG9iago0IDAgb2JqCjw8Ci9Db250ZW50cyA4IDAgUiAvTWVkaWFCb3ggWyAwIDAgNjEyIDc5MiBdIC9QYXJl"
    "bnQgNyAwIFIgL1Jlc291cmNlcyA8PAovRm9udCAxIDAgUiAvUHJvY1NldCBbIC9QREYgL1RleHQgL0ltYWdlQiAv"
    "SW1hZ2VDIC9JbWFnZUkgXSAvWE9iamVjdCA8PAovRm9ybVhvYi4wNDUwMjk3MmRlMzM2NmMxMGVkODgyZmU5NDM5"
    "ZjgxNSAzIDAgUgo+Pgo+PiAvUm90YXRlIDAgL1RyYW5zIDw8Cgo+PiAKICAvVHlwZSAvUGFnZQo+PgplbmRvYmoK"
    "NSAwIG9iago8PAovUGFnZU1vZGUgL1VzZU5vbmUgL1BhZ2VzIDcgMCBSIC9UeXBlIC9DYXRhbG9nCj4+CmVuZG9i"
    "ago2IDAgb2JqCjw8Ci9BdXRob3IgKGFub255bW91cykgL0NyZWF0aW9uRGF0ZSAoRDoyMDI2MDYzMDEwMjE0Mysw"
    "MicwMCcpIC9DcmVhdG9yIChhbm9ueW1vdXMpIC9LZXl3b3JkcyAoKSAvTW9kRGF0ZSAoRDoyMDI2MDYzMDEwMjE0"
    "MyswMicwMCcpIC9Qcm9kdWNlciAoUmVwb3J0TGFiIFBERiBMaWJyYXJ5IC0gXChvcGVuc291cmNlXCkpIAogIC9T"
    "dWJqZWN0ICh1bnNwZWNpZmllZCkgL1RpdGxlICh1bnRpdGxlZCkgL1RyYXBwZWQgL0ZhbHNlCj4+CmVuZG9iago3"
    "IDAgb2JqCjw8Ci9Db3VudCAxIC9LaWRzIFsgNCAwIFIgXSAvVHlwZSAvUGFnZXMKPj4KZW5kb2JqCjggMCBvYmoK"
    "PDwKL0ZpbHRlciBbIC9BU0NJSTg1RGVjb2RlIC9GbGF0ZURlY29kZSBdIC9MZW5ndGggMTcxCj4+CnN0cmVhbQpH"
    "YXEzXDN0Jm5aJ1NaO1pNRE04IkhaLEBdYVhaI2Q3Q1t1TlZNUz1aNXEyJEQyYTBMUEBsU2prNXQxXWNUN0Y+Qipa"
    "MXQlJlNycCFmalFEMm0vMDR0PnRnPzVLblNXSmEjJS1XVjtqUWUrXWxiTWMvMDNLSTFQdDI6VD5LLDRXWCptLDFE"
    "XXVdY0tQaUZJZCVPVVU8LGVCLFZSPFwoRzYhPSY6UTxINm9Ofj5lbmRzdHJlYW0KZW5kb2JqCnhyZWYKMCA5CjAw"
    "MDAwMDAwMDAgNjU1MzUgZiAKMDAwMDAwMDA2MSAwMDAwMCBuIAowMDAwMDAwMDkyIDAwMDAwIG4gCjAwMDAwMDAx"
    "OTkgMDAwMDAgbiAKMDAwMDAwMDQwNCAwMDAwMCBuIAowMDAwMDAwNjYwIDAwMDAwIG4gCjAwMDAwMDA3MjggMDAw"
    "MDAgbiAKMDAwMDAwMDk4OSAwMDAwMCBuIAowMDAwMDAxMDQ4IDAwMDAwIG4gCnRyYWlsZXIKPDwKL0lEIApbPDE0"
    "ODFjMjJlYmMxZDRmMzU0ZTQ5MmZjMmFkMzdlMjhmPjwxNDgxYzIyZWJjMWQ0ZjM1NGU0OTJmYzJhZDM3ZTI4Zj5d"
    "CiUgUmVwb3J0TGFiIGdlbmVyYXRlZCBQREYgZG9jdW1lbnQgLS0gZGlnZXN0IChvcGVuc291cmNlKQoKL0luZm8g"
    "NiAwIFIKL1Jvb3QgNSAwIFIKL1NpemUgOQo+PgpzdGFydHhyZWYKMTMwOQolJUVPRgo="
)


def _make_png_bytes(*, color: tuple[int, int, int] = (255, 0, 0)) -> bytes:
    """Create a tiny PNG suitable for embedding in PDF/PPTX fixtures."""
    img = Image.new("RGB", (4, 4), color)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _make_docx(path: Path, paragraphs: list[tuple[str, str]]) -> None:
    """Create DOCX with (style, text) pairs. Style '' means normal."""
    doc = DocxDocument()
    for style, text in paragraphs:
        doc.add_paragraph(text, style=style if style else None)
    doc.save(str(path))


def _make_pptx(path: Path, slides_text: list[list[str]], *, embed_image: bool = False) -> None:
    """Create PPTX with text boxes per slide, optionally embedding an image."""
    prs = Presentation()
    for texts in slides_text:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        for i, text in enumerate(texts):
            left = Inches(1)
            top = Inches(1 + i)
            txBox = slide.shapes.add_textbox(left, top, Inches(5), Inches(0.5))
            txBox.text_frame.text = text
        if embed_image:
            img_bytes = _make_png_bytes(color=(0, 0, 255))
            with tempfile.NamedTemporaryFile(suffix=".png", dir=path.parent, delete=False) as tmp:
                tmp.write(img_bytes)
                img_tmp_path = Path(tmp.name)
            try:
                slide.shapes.add_picture(str(img_tmp_path), Inches(3), Inches(3), Inches(1), Inches(1))
            finally:
                img_tmp_path.unlink(missing_ok=True)
    prs.save(str(path))


def _make_xlsx(path: Path, sheets: dict[str, list[list[str]]]) -> None:
    wb = Workbook()
    first = True
    for name, rows in sheets.items():
        ws = wb.active if first else wb.create_sheet(title=name)
        if first:
            ws.title = name
            first = False
        for row in rows:
            ws.append(row)
    wb.save(str(path))


def _make_csv(path: Path, rows: list[list[str]]) -> None:
    with open(path, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerows(rows)


# ---------------------------------------------------------------------------
# is_image_file / is_supported_document
# ---------------------------------------------------------------------------

class TestFileTypeDetection:
    @pytest.mark.parametrize("ext,expected", [
        (".png", True), (".jpg", True), (".jpeg", True), (".gif", True),
        (".webp", True), (".svg", True), (".bmp", True),
        (".pdf", False), (".docx", False), (".txt", False),
    ])
    def test_is_image_file(self, ext: str, expected: bool) -> None:
        assert is_image_file(Path(f"test{ext}")) == expected

    @pytest.mark.parametrize("ext,expected", [
        (".pdf", True), (".docx", True), (".pptx", True),
        (".xlsx", True), (".xls", True), (".csv", True),
        (".png", False), (".txt", False), (".md", False),
    ])
    def test_is_supported_document(self, ext: str, expected: bool) -> None:
        assert is_supported_document(Path(f"test{ext}")) == expected


# ---------------------------------------------------------------------------
# _truncate
# ---------------------------------------------------------------------------

class TestTruncate:
    def test_short_text_not_truncated(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """Scenario: document under truncation limit → no truncation warning."""
        import services.document_reader as mod
        monkeypatch.setattr(mod, "MAX_TEXT_LENGTH", 10_000)
        p = tmp_path / "short.csv"
        _make_csv(p, [["hello", "world"]])
        result = read_document(p)
        assert not any("truncated" in w.lower() for w in result.warnings)
        assert "truncated" not in result.markdown

    def test_long_text_truncated(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """Scenario: document over truncation limit → text cut + warning added."""
        import services.document_reader as mod
        monkeypatch.setattr(mod, "MAX_TEXT_LENGTH", 20)
        p = tmp_path / "long.csv"
        _make_csv(p, [["x" * 50]])
        result = read_document(p)
        assert any("truncated" in w.lower() for w in result.warnings)
        assert "truncated" in result.markdown


# ---------------------------------------------------------------------------
# Error cases
# ---------------------------------------------------------------------------

class TestReadDocumentErrors:
    def test_file_not_found(self, tmp_path: Path) -> None:
        with pytest.raises(FileNotFoundError):
            read_document(tmp_path / "nope.pdf")

    def test_unsupported_format(self, tmp_path: Path) -> None:
        p = tmp_path / "test.xyz"
        p.write_text("data")
        with pytest.raises(ValueError, match="Unsupported format"):
            read_document(p)

    def test_file_too_large(self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
        """Scenario: file exceeds MAX_FILE_SIZE_MB → ValueError raised."""
        import services.document_reader as mod
        monkeypatch.setattr(mod, "MAX_FILE_SIZE_MB", 0.0001)  # ~100 bytes
        p = tmp_path / "big.pdf"
        p.write_bytes(b"x" * 200)
        with pytest.raises(ValueError, match="too large"):
            read_document(p)


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

class TestReadPdf:
    def test_basic_text_extraction(self, tmp_path: Path) -> None:
        """Scenario: 2-page PDF with known text → markdown contains both pages."""
        p = tmp_path / "test.pdf"
        p.write_bytes(_TWO_PAGE_TEXT_PDF)

        result = read_document(p)

        assert isinstance(result, DocumentContent)
        assert result.page_count == 2
        assert "Page one content" in result.markdown
        assert "Page two content" in result.markdown
        assert "## Page 1" in result.markdown
        assert "## Page 2" in result.markdown

    def test_embedded_image_extraction(self, tmp_path: Path) -> None:
        """Scenario: PDF with an embedded image → image extracted as base64."""
        p = tmp_path / "img.pdf"
        p.write_bytes(_PDF_WITH_EMBEDDED_IMAGE)

        result = read_document(p)

        assert len(result.images) >= 1
        img = result.images[0]
        assert isinstance(img, EmbeddedImage)
        assert img.mime_type == "image/png"
        assert len(img.data_b64) > 0
        assert img.width is not None
        assert img.height is not None


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

class TestReadDocx:
    def test_paragraphs_and_headings(self, tmp_path: Path) -> None:
        """Scenario: DOCX with heading + normal para → markdown uses # syntax."""
        p = tmp_path / "test.docx"
        _make_docx(p, [
            ("Heading 1", "Main Title"),
            ("", "Normal paragraph text."),
            ("Heading 2", "Subtitle"),
        ])

        result = read_document(p)

        assert "# Main Title" in result.markdown
        assert "Normal paragraph text." in result.markdown
        assert "## Subtitle" in result.markdown

    def test_table_extraction(self, tmp_path: Path) -> None:
        """Scenario: DOCX with a table → markdown table rendered."""
        doc = DocxDocument()
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "A"
        table.cell(0, 1).text = "B"
        table.cell(1, 0).text = "C"
        table.cell(1, 1).text = "D"
        p = tmp_path / "table.docx"
        doc.save(str(p))

        result = read_document(p)

        assert "| A | B |" in result.markdown
        assert "| C | D |" in result.markdown
        assert "---" in result.markdown

    def test_table_and_paragraph_ordering(self, tmp_path: Path) -> None:
        """Scenario: table sandwiched between two paragraphs → table appears in document order."""
        doc = DocxDocument()
        doc.add_paragraph("Before table")
        tbl = doc.add_table(rows=1, cols=2)
        tbl.cell(0, 0).text = "Col1"
        tbl.cell(0, 1).text = "Col2"
        doc.add_paragraph("After table")
        p = tmp_path / "order.docx"
        doc.save(str(p))

        result = read_document(p)

        before_pos = result.markdown.index("Before table")
        table_pos = result.markdown.index("| Col1 |")
        after_pos = result.markdown.index("After table")
        assert before_pos < table_pos < after_pos


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

class TestReadPptx:
    def test_slide_text_extraction(self, tmp_path: Path) -> None:
        """Scenario: PPTX with 2 slides → markdown has slide headings and text."""
        p = tmp_path / "test.pptx"
        _make_pptx(p, [["Slide one text"], ["Slide two", "More text"]])

        result = read_document(p)

        assert result.page_count == 2
        assert "## Slide 1" in result.markdown
        assert "Slide one text" in result.markdown
        assert "## Slide 2" in result.markdown
        assert "More text" in result.markdown

    def test_embedded_image_extraction(self, tmp_path: Path) -> None:
        """Scenario: PPTX with embedded image → image extracted as base64."""
        p = tmp_path / "img.pptx"
        _make_pptx(p, [["Slide with image"]], embed_image=True)

        result = read_document(p)

        assert len(result.images) >= 1
        assert result.images[0].mime_type.startswith("image/")
        assert "Slide 1" in result.images[0].label


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------

class TestReadXlsx:
    def test_single_sheet(self, tmp_path: Path) -> None:
        """Scenario: XLSX with one sheet → markdown table rendered."""
        p = tmp_path / "test.xlsx"
        _make_xlsx(p, {"Data": [["Name", "Value"], ["Alice", "10"], ["Bob", "20"]]})

        result = read_document(p)

        assert "## Sheet: Data" in result.markdown
        assert "| Name | Value |" in result.markdown
        assert "| Alice | 10 |" in result.markdown

    def test_multiple_sheets(self, tmp_path: Path) -> None:
        """Scenario: XLSX with 2 sheets → both appear in output."""
        p = tmp_path / "multi.xlsx"
        _make_xlsx(p, {
            "Sheet1": [["A", "B"]],
            "Sheet2": [["X", "Y"]],
        })

        result = read_document(p)

        assert "## Sheet: Sheet1" in result.markdown
        assert "## Sheet: Sheet2" in result.markdown


# ---------------------------------------------------------------------------
# CSV
# ---------------------------------------------------------------------------

class TestReadCsv:
    def test_basic_csv(self, tmp_path: Path) -> None:
        """Scenario: CSV with header + rows → markdown table."""
        p = tmp_path / "test.csv"
        _make_csv(p, [["Name", "Age"], ["Alice", "30"], ["Bob", "25"]])

        result = read_document(p)

        assert "| Name | Age |" in result.markdown
        assert "| Alice | 30 |" in result.markdown
        assert "---" in result.markdown

    def test_pipe_in_cells_escaped(self, tmp_path: Path) -> None:
        """Scenario: CSV cell contains pipe character → escaped in markdown."""
        p = tmp_path / "pipe.csv"
        _make_csv(p, [["Col"], ["val|ue"]])

        result = read_document(p)

        assert "\\|" in result.markdown
