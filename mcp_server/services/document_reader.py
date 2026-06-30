"""Extract text and embedded images from common document formats.

Returns structured markdown + base64 image references so MCP clients
(Cursor, Claude Code, Copilot, Gemini CLI, Cline, etc.) can consume
document content regardless of their native file-format support.

Supported formats:
  - PDF  (.pdf)        — via pypdf (BSD-3-Clause)
  - DOCX (.docx)       — via python-docx
  - PPTX (.pptx)       — via python-pptx
  - XLSX (.xlsx, .xls) — via openpyxl
  - CSV  (.csv)        — via stdlib csv

PDF support uses pypdf (BSD-3-Clause) so specflow remains MIT-compatible
and avoids AGPL obligations from PyMuPDF/MuPDF (DEP-001).

Images embedded in PDF / PPTX are extracted as base64 and returned
alongside the markdown so the IDE's vision model can interpret them.
No local OCR is performed.
"""

from __future__ import annotations

import base64
import csv
import io
import logging
from collections.abc import Callable
from dataclasses import dataclass, field
from pathlib import Path

from PIL import Image
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader

logger = logging.getLogger(__name__)

MAX_FILE_SIZE_MB = 50
MAX_CSV_ROWS = 5_000
MAX_TEXT_LENGTH = 400_000  # ~400 KB of text output
MAX_IMAGES = 50

IMAGE_EXTENSIONS = frozenset({".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg", ".bmp"})
SUPPORTED_EXTENSIONS = frozenset({".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".csv"})

_IMG_MIME = {
    "png": "image/png",
    "jpeg": "image/jpeg",
    "jpg": "image/jpeg",
    "webp": "image/webp",
    "gif": "image/gif",
    "bmp": "image/bmp",
    "tiff": "image/tiff",
    "svg": "image/svg+xml",
}


@dataclass
class EmbeddedImage:
    """An image extracted from a document."""

    data_b64: str
    mime_type: str
    label: str
    width: int | None = None
    height: int | None = None


@dataclass
class DocumentContent:
    """Structured result from document extraction."""

    markdown: str
    images: list[EmbeddedImage] = field(default_factory=list)
    page_count: int | None = None
    warnings: list[str] = field(default_factory=list)


def _truncate(text: str) -> tuple[str, bool]:
    if len(text) <= MAX_TEXT_LENGTH:
        return text, False
    return text[:MAX_TEXT_LENGTH] + "\n\n[… content truncated at 400 KB …]", True


def _pypdf_image_to_png(image_file_object: object) -> tuple[bytes, int, int]:
    """Decode a pypdf image object to PNG bytes and pixel dimensions."""
    pil_image: Image.Image | None = None
    try:
        pil_image = image_file_object.image  # type: ignore[attr-defined]
    except Exception:
        pass
    if pil_image is None:
        try:
            pil_image = image_file_object.decode_as_image()  # type: ignore[attr-defined]
        except Exception:
            raw = image_file_object.data  # type: ignore[attr-defined]
            pil_image = Image.open(io.BytesIO(raw))
    if pil_image.mode not in ("RGB", "RGBA", "L", "LA"):
        pil_image = pil_image.convert("RGBA")
    buf = io.BytesIO()
    pil_image.save(buf, format="PNG")
    return buf.getvalue(), pil_image.width, pil_image.height


# ---------------------------------------------------------------------------
# PDF via pypdf
# ---------------------------------------------------------------------------

def _read_pdf(path: Path) -> DocumentContent:
    parts: list[str] = []
    images: list[EmbeddedImage] = []
    warnings: list[str] = []
    image_cap_reached = False

    reader = PdfReader(str(path))
    for page_num, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        parts.append(f"## Page {page_num}\n\n{text.strip()}")

        if not image_cap_reached:
            for img_index, name in enumerate(page.images.keys(), 1):
                if len(images) >= MAX_IMAGES:
                    warnings.append(f"Image extraction capped at {MAX_IMAGES} images.")
                    image_cap_reached = True
                    break
                try:
                    image_file = page.images[name]
                    img_bytes, width, height = _pypdf_image_to_png(image_file)
                    images.append(EmbeddedImage(
                        data_b64=base64.b64encode(img_bytes).decode("ascii"),
                        mime_type="image/png",
                        label=f"Page {page_num} — image {img_index}",
                        width=width,
                        height=height,
                    ))
                except Exception as exc:
                    warnings.append(
                        f"Page {page_num}, image {img_index}: extraction failed ({exc})"
                    )

    page_count = len(reader.pages)
    md, truncated = _truncate("\n\n".join(parts))
    if truncated:
        warnings.append("Text content was truncated (exceeded 400 KB).")
    return DocumentContent(markdown=md, images=images, page_count=page_count, warnings=warnings)


# ---------------------------------------------------------------------------
# DOCX via python-docx
# ---------------------------------------------------------------------------

def _read_docx(path: Path) -> DocumentContent:
    doc = DocxDocument(str(path))
    parts: list[str] = []
    images: list[EmbeddedImage] = []
    warnings: list[str] = []

    # Build lookup maps keyed by XML element to preserve document order
    # when iterating body children (paragraphs and tables are interleaved).
    para_by_elem = {p._element: p for p in doc.paragraphs}
    table_by_elem = {t._element: t for t in doc.tables}

    for child in doc.element.body:
        if child in para_by_elem:
            para = para_by_elem[child]
            style = para.style.name if para.style else ""
            text = para.text.strip()
            if not text:
                continue
            if style.startswith("Heading"):
                level = style.replace("Heading", "").strip()
                try:
                    hashes = "#" * int(level)
                except ValueError:
                    hashes = "##"
                parts.append(f"{hashes} {text}")
            else:
                parts.append(text)
        elif child in table_by_elem:
            table = table_by_elem[child]
            rows_md: list[str] = []
            for i, row in enumerate(table.rows):
                cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
                rows_md.append("| " + " | ".join(cells) + " |")
                if i == 0:
                    rows_md.append("| " + " | ".join(["---"] * len(cells)) + " |")
            parts.append("\n".join(rows_md))

    img_index = 0
    for rel in doc.part.rels.values():
        if "image" in rel.reltype and not rel.is_external:
            img_index += 1
            if len(images) >= MAX_IMAGES:
                warnings.append(f"Image extraction capped at {MAX_IMAGES} images.")
                break
            try:
                blob = rel.target_part.blob
                ext = Path(rel.target_part.partname).suffix.lstrip(".").lower()
                mime = _IMG_MIME.get(ext, f"image/{ext}")
                images.append(EmbeddedImage(
                    data_b64=base64.b64encode(blob).decode("ascii"),
                    mime_type=mime,
                    label=f"Document image {img_index}",
                ))
            except Exception as exc:
                warnings.append(f"Image {img_index}: extraction failed ({exc})")

    md, truncated = _truncate("\n\n".join(parts))
    if truncated:
        warnings.append("Text content was truncated (exceeded 400 KB).")
    return DocumentContent(markdown=md, images=images, warnings=warnings)


# ---------------------------------------------------------------------------
# PPTX via python-pptx
# ---------------------------------------------------------------------------

def _read_pptx(path: Path) -> DocumentContent:
    prs = Presentation(str(path))
    parts: list[str] = []
    images: list[EmbeddedImage] = []
    warnings: list[str] = []
    image_cap_reached = False

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = [f"## Slide {slide_num}"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_parts.append(text)

            if shape.has_table:
                rows_md: list[str] = []
                for i, row in enumerate(shape.table.rows):
                    cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
                    rows_md.append("| " + " | ".join(cells) + " |")
                    if i == 0:
                        rows_md.append("| " + " | ".join(["---"] * len(cells)) + " |")
                slide_parts.append("\n".join(rows_md))

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and not image_cap_reached:
                try:
                    blob = shape.image.blob
                    # width/height are omitted: shape dimensions are in EMU
                    # (not pixels) and not meaningful to the vision model.
                    images.append(EmbeddedImage(
                        data_b64=base64.b64encode(blob).decode("ascii"),
                        mime_type=shape.image.content_type,
                        label=f"Slide {slide_num} — image",
                    ))
                    if len(images) >= MAX_IMAGES:
                        warnings.append(f"Image extraction capped at {MAX_IMAGES} images.")
                        image_cap_reached = True
                except Exception as exc:
                    warnings.append(f"Slide {slide_num}: image extraction failed ({exc})")

        # Read speaker notes from the notes body placeholder (idx=1) only,
        # skipping the slide-image thumbnail placeholder (idx=0).
        if slide.has_notes_slide:
            for shape in slide.notes_slide.shapes:
                if (
                    shape.has_text_frame
                    and shape.placeholder_format is not None
                    and shape.placeholder_format.idx == 1
                ):
                    notes = shape.text_frame.text.strip()
                    if notes:
                        slide_parts.append(f"\n> **Speaker notes:** {notes}")
                    break

        parts.append("\n\n".join(slide_parts))

    md, truncated = _truncate("\n\n---\n\n".join(parts))
    if truncated:
        warnings.append("Text content was truncated (exceeded 400 KB).")
    return DocumentContent(
        markdown=md, images=images, page_count=len(prs.slides), warnings=warnings,
    )


# ---------------------------------------------------------------------------
# XLSX via openpyxl
# ---------------------------------------------------------------------------

def _read_xlsx(path: Path) -> DocumentContent:
    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    warnings: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"## Sheet: {sheet_name}")

        rows_md: list[str] = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            cells = [str(cell if cell is not None else "") for cell in row]
            cells = [c.replace("|", "\\|") for c in cells]
            rows_md.append("| " + " | ".join(cells) + " |")
            if i == 0:
                rows_md.append("| " + " | ".join(["---"] * len(cells)) + " |")
            if i >= MAX_CSV_ROWS:
                warnings.append(f"Sheet '{sheet_name}': truncated at {MAX_CSV_ROWS} rows.")
                break
        parts.append("\n".join(rows_md))

    wb.close()
    md, truncated = _truncate("\n\n".join(parts))
    if truncated:
        warnings.append("Text content was truncated (exceeded 400 KB).")
    return DocumentContent(markdown=md, warnings=warnings)


# ---------------------------------------------------------------------------
# CSV via stdlib
# ---------------------------------------------------------------------------

def _read_csv(path: Path) -> DocumentContent:
    warnings: list[str] = []
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        sniffer = csv.Sniffer()
        sample = f.read(8192)
        f.seek(0)
        try:
            dialect = sniffer.sniff(sample)
        except csv.Error:
            dialect = csv.excel
        reader = csv.reader(f, dialect)

        rows_md: list[str] = []
        for i, row in enumerate(reader):
            cells = [c.replace("|", "\\|") for c in row]
            rows_md.append("| " + " | ".join(cells) + " |")
            if i == 0:
                rows_md.append("| " + " | ".join(["---"] * len(cells)) + " |")
            if i >= MAX_CSV_ROWS:
                warnings.append(f"CSV truncated at {MAX_CSV_ROWS} rows.")
                break

    md, truncated = _truncate("\n".join(rows_md))
    if truncated:
        warnings.append("Text content was truncated (exceeded 400 KB).")
    return DocumentContent(markdown=md, warnings=warnings)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

_READERS: dict[str, Callable[[Path], DocumentContent]] = {
    ".pdf": _read_pdf,
    ".docx": _read_docx,
    ".pptx": _read_pptx,
    ".xlsx": _read_xlsx,
    ".xls": _read_xlsx,
    ".csv": _read_csv,
}


def is_image_file(path: Path) -> bool:
    return path.suffix.lower() in IMAGE_EXTENSIONS


def is_supported_document(path: Path) -> bool:
    return path.suffix.lower() in SUPPORTED_EXTENSIONS


def read_document(path: Path) -> DocumentContent:
    """Read a document and return structured markdown + embedded images.

    Raises:
        FileNotFoundError: if the file does not exist.
        ValueError: if the format is unsupported or the file is too large.
    """
    if not path.exists():
        raise FileNotFoundError(f"File not found: {path}")

    size_mb = path.stat().st_size / (1024 * 1024)
    if size_mb > MAX_FILE_SIZE_MB:
        raise ValueError(
            f"File too large ({size_mb:.1f} MB). Maximum supported size is {MAX_FILE_SIZE_MB} MB."
        )

    ext = path.suffix.lower()
    reader = _READERS.get(ext)
    if reader is None:
        raise ValueError(
            f"Unsupported format: {ext}. "
            f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    return reader(path)
